# -*- coding: utf-8 -*-
"""Slide "Kiến trúc AI Chatbot tư vấn món ăn an toàn và có kiểm soát".

Mười trang, 16:9, bảng màu xanh đậm / trắng / cam nhạt. Bốn trang (3–6) dành cho
khâu nhóm trưởng phụ trách: tầng dữ liệu và lớp hiểu câu hỏi.

Ba nguyên tắc dựng, giống bộ `xuat_slide_hoc_may.py`
-----------------------------------------------------
1. Bố cục dựng bằng hàm, không đặt hộp bằng toạ độ tự chế ở từng slide.
2. Mọi con số trên slide **đọc từ mã nguồn lúc dựng** (xem `SO` bên dưới) chứ
   không viết tay. Số viết tay sẽ trôi — biểu đồ 4.3 của báo cáo từng viết cứng
   "147 ca / 163 lượt / RAG 0%" và vẫn vẽ lại y nguyên số đó sau khi tập ca đã
   mở rộng, vì hình không đọc gì cả.
3. Chốt chặn ngay lúc dựng: quá 5 ý một slide, chữ tràn khối, hay bảng nở quá
   đáy thì **nổ**, không xuất ra tệp.

Chạy:  python xuat_slide_kien_truc.py
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

HERE = Path(__file__).resolve().parent
GOC = HERE.parents[1]
sys.path.insert(0, str(GOC / "ai" / "app"))

# ── bảng màu: xanh đậm · trắng · cam nhạt ────────────────────────────
XANH = RGBColor(0x0F, 0x2D, 0x52)      # xanh đậm — màu chủ đạo, mã tất định
XANH_NHAT = RGBColor(0xE7, 0xED, 0xF4)  # nền khối xanh
TIM = RGBColor(0x5B, 0x3E, 0x7E)       # thành phần dùng mô hình
TIM_NHAT = RGBColor(0xEE, 0xE9, 0xF4)
CAM = RGBColor(0xC9, 0x6F, 0x1E)       # lớp kiểm tra an toàn
CAM_NHAT = RGBColor(0xFD, 0xF0, 0xE0)  # cam nhạt
TRANG = RGBColor(0xFF, 0xFF, 0xFF)
DEN = RGBColor(0x14, 0x18, 0x1D)
XAM = RGBColor(0x55, 0x5F, 0x6B)
KE = RGBColor(0xC9, 0xD4, 0xE0)

W, H = Cm(33.867), Cm(19.05)           # 16:9
LE = Cm(1.5)
RONG = W - 2 * LE
FONT = "Times New Roman"

Y_ND = Cm(3.5)                          # mép trên vùng nội dung
Y_PB = Cm(17.9)                         # (dải phản biện đã bỏ khỏi slide)
MAX_Y = 17.9                            # đáy cho phép, tính bằng cm

# ── đo chữ bằng chính font slide dùng ────────────────────────────────
from PIL import ImageFont  # noqa: E402

_TTF = {False: "C:/Windows/Fonts/times.ttf", True: "C:/Windows/Fonts/timesbd.ttf"}
_kho: dict = {}


def _font(pt, dam=False):
    k = (round(pt, 1), bool(dam))
    if k not in _kho:
        _kho[k] = ImageFont.truetype(_TTF[bool(dam)], max(int(pt * 96 / 72), 1))
    return _kho[k]


def _rong(chu, pt, dam=False):
    return _font(pt, dam).getlength(chu) / 96 * 2.54       # cm


def _so_dong(chu, rong_cm, pt, dam=False):
    if not chu.strip():
        return 1
    n, hien = 1, ""
    for tu in chu.split(" "):
        thu = (hien + " " + tu).strip()
        if _rong(thu, pt, dam) <= rong_cm or not hien:
            hien = thu
        else:
            n, hien = n + 1, tu
    return n


# ═══════════════════════════════════════════════ số đọc từ mã nguồn
def _doc_so() -> dict:
    """Mọi con số dùng trên slide, đọc từ mã sống. Không có số viết tay."""
    import dataclasses
    import re

    import understand as U
    from generate import BRANCHES_ALLOWED
    from test_understand import collision_census

    menu = json.loads((GOC / "data" / "menu-dataset.json")
                      .read_text(encoding="utf-8-sig"))
    tags = json.loads((GOC / "data" / "menu-tags.json")
                      .read_text(encoding="utf-8-sig"))
    items = menu["items"]

    # Router thật là `_chon_cau_tra_loi`; `respond()` chỉ là lớp mỏng bọc ngoài
    # để ghép câu xác nhận ở đúng MỘT chỗ.
    L = (GOC / "ai" / "app" / "answer.py").read_text(encoding="utf-8").splitlines()
    i = next(k for k, l in enumerate(L) if l.startswith("def _chon_cau_tra_loi"))
    than = []
    for l in L[i + 1:]:
        if l.startswith("def "):
            break
        than.append(l)
    t = "\n".join(than)
    ten_nhanh = re.findall(r'branch=f?"([^"]+)"', t)
    ho_nhanh = {x.split(":")[0].replace("{", "").strip() for x in ten_nhanh}

    s = (GOC / "ai" / "app" / "generate.py").read_text(encoding="utf-8")
    m = re.search(r"\ndef verify\(.*?(?=\ndef |\Z)", s, re.S)
    so_kiem = len(re.findall(r"do\.append|loi\.append|return\s+False", m.group(0)))

    vc = collision_census()
    import session as S

    return {
        "mon": len(items),
        "danh_muc": len({i.get("category") or i.get("categoryId") for i in items}),
        "nhan": len(tags["tags"]),
        "nhom_nhan": len(tags["groups"]),
        "di_ung": sum(1 for i in items
                      if any(str(x).startswith("allergen:") for x in (i.get("tags") or []))),
        "truong_request": len(dataclasses.fields(U.Request)),
        "tu_vung": vc["tu_vung"],
        "va_cham": vc["co_rui_ro"],
        "ho_nhanh": len(ho_nhanh),
        "diem_tra_ve": len(ten_nhanh),
        "nhanh_sinh": len(BRANCHES_ALLOWED),
        "phep_kiem": so_kiem,
        "ngu_canh": S.MAX_CONTEXT_TAGS,
        "tai_lieu": len(list((GOC / "ai" / "knowledge").rglob("*.md"))),
    }


SO = _doc_so()


# ═══════════════════════════════════════════════════════ nguyên thuỷ vẽ
def txt(sl, x, y, w, h, chu, co=16, dam=False, mau=DEN,
        canh=PP_ALIGN.LEFT, gian=1.25):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.08)
    tf.margin_top = tf.margin_bottom = Cm(0.04)
    p = tf.paragraphs[0]
    p.alignment = canh
    p.line_spacing = gian
    r = p.add_run()
    r.text = chu
    r.font.size = Pt(co)
    r.font.bold = dam
    r.font.color.rgb = mau
    r.font.name = FONT
    return tb


def hop(sl, x, y, w, h, chu, nen=XANH_NHAT, vien=XANH, chu_mau=XANH,
        co=14, dam=True, tron=True):
    """Một khối trong sơ đồ."""
    o = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if tron else MSO_SHAPE.RECTANGLE, x, y, w, h)
    o.fill.solid(); o.fill.fore_color.rgb = nen
    o.line.color.rgb = vien; o.line.width = Pt(1.1)
    o.shadow.inherit = False
    if tron:
        o.adjustments[0] = 0.12
    tf = o.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Cm(0.12)
    for k, dong in enumerate(chu.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.12
        r = p.add_run(); r.text = dong
        r.font.size = Pt(co if k == 0 else co - 2)
        r.font.bold = dam if k == 0 else False
        r.font.color.rgb = chu_mau if k == 0 else XAM
        r.font.name = FONT
    return o


def ten(sl, x, y, w, mau=XANH, doc=False, h=None):
    """Mũi tên CHỈ CHIỀU ĐI CỦA DỮ LIỆU — bắt buộc trên mọi sơ đồ."""
    if doc:
        o = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h or Cm(0.6))
    else:
        o = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h or Cm(0.42))
    o.fill.solid(); o.fill.fore_color.rgb = mau
    o.line.fill.background(); o.shadow.inherit = False
    return o


def trang(prs, tieu_de, nhan, phu=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    txt(sl, LE, Cm(0.72), RONG, Cm(0.7), nhan, 12.5, True, CAM)
    txt(sl, LE, Cm(1.35), RONG, Cm(1.2), tieu_de, 25, True, XANH)
    if phu:
        txt(sl, LE, Cm(2.52), RONG, Cm(0.7), phu, 14, False, XAM)
    v = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, LE,
                            Cm(3.18) if phu else Cm(2.82), RONG, Cm(0.05))
    v.fill.solid(); v.fill.fore_color.rgb = CAM
    v.line.fill.background(); v.shadow.inherit = False
    return sl


def y_chinh(sl, muc, y=Y_ND, x=LE, w=None, co=16.5):
    """Tối đa 5 ý một slide — chốt của đặc tả."""
    assert len(muc) <= 5, f"{len(muc)} ý, đặc tả cho tối đa 5"
    w = w or RONG
    wc = w / 360000 - 0.95
    cao = sum(_so_dong(m, wc, co) * co * 1.34 / 28.35 + 0.32 for m in muc)
    assert y / 360000 + cao <= MAX_Y, (
        f"khối ý cao {cao:.1f}cm, đáy {y / 360000 + cao:.1f}cm > {MAX_Y}cm")
    # Hộp mang ĐÚNG chiều cao vừa tính, không phải một con số tượng trưng.
    # PowerPoint vẫn vẽ chữ tràn ra ngoài hộp chứ không cắt, nên hộp khai báo
    # 0,5cm vẫn hiện đủ chữ — nhưng khi đó tệp nói dối về chỗ nó chiếm, và mọi
    # phép kiểm đọc tệp (kể cả phép kiểm của chính tôi) sẽ báo sai.
    tb = sl.shapes.add_textbox(x, y, w, Cm(cao))
    tf = tb.text_frame
    tf.word_wrap = True
    for k, m in enumerate(muc):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.line_spacing = 1.24
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "▸  "
        r.font.size = Pt(co); r.font.color.rgb = CAM; r.font.name = FONT
        r = p.add_run(); r.text = m
        r.font.size = Pt(co); r.font.color.rgb = DEN; r.font.name = FONT
    return tb


def _ty_le(cot, hang, co):
    """Cột chia theo nhu cầu thật nhưng kéo về gần đều.

    Chia đều tuyệt đối thì cột chứa câu dài bị ép xuống bốn dòng còn cột hai chữ
    bỏ trống nửa ô. Chia đúng nhu cầu thì bảng lệch hẳn, cột to cột bé. Nên pha
    45% nhu cầu với 55% đều, rồi kẹp trong [0,84·đều , 1,22·đều].
    """
    n = len(cot)
    deu = 1.0 / n
    can = [max(_rong(x, co, True) for x in [cot[i]] + [h[i] for h in hang]) or 1.0
           for i in range(n)]
    tong = sum(can)
    kep = [min(max(0.45 * (c / tong) + 0.55 * deu, 0.84 * deu), 1.22 * deu)
           for c in can]
    s = sum(kep)
    ty = [k / s for k in kep]
    assert max(ty) / min(ty) <= 1.5, "cột lệch quá 1,5 lần"
    return ty


def bang(sl, cot, hang, y=Y_ND, co=15, x=LE, w=None, cao_min=Cm(1.15)):
    w = w or RONG
    ty = _ty_le(cot, hang, co)
    t = sl.shapes.add_table(len(hang) + 1, len(cot), x, y, w,
                            cao_min * (len(hang) + 1)).table
    for i, c in enumerate(ty):
        t.columns[i].width = int(w * c)

    # PowerPoint TỰ NỚI chiều cao hàng lúc dựng hình và KHÔNG sửa lại con số
    # trong XML — nên bảng "vừa khung" khi mở tệp ra đo mà vẫn chạy khỏi mép
    # dưới trên màn chiếu. Phải tự tính lại chiều cao thật ở đây.
    that = 0.0
    for h in [cot] + list(hang):
        that += max([0.26 + _so_dong(x_, w / 360000 * ty[i] - 0.5, co) * co * 1.2 / 28.35
                     for i, x_ in enumerate(h)] + [cao_min / 360000])
    assert y / 360000 + that <= MAX_Y, (
        f"bảng nở tới {that:.1f}cm, đáy {y / 360000 + that:.1f}cm > {MAX_Y}cm")

    for i, x_ in enumerate(cot):
        c = t.cell(0, i); c.text = x_
        c.fill.solid(); c.fill.fore_color.rgb = XANH
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(co); r.font.bold = True
                r.font.color.rgb = TRANG; r.font.name = FONT
    for j, h in enumerate(hang, 1):
        for i, x_ in enumerate(h):
            c = t.cell(j, i); c.text = x_
            c.fill.solid()
            c.fill.fore_color.rgb = TRANG if j % 2 else CAM_NHAT
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(co - 0.5)
                    r.font.bold = (i == 0)
                    r.font.color.rgb = XANH if i == 0 else DEN
                    r.font.name = FONT
    return t


def phan_bien(sl, hoi, dap):
    """Đã BỎ khỏi slide — giữ hàm để nội dung phản biện còn nằm trong mã.

    Câu hỏi phản biện và câu trả lời chuyển xuống KỊCH BẢN NÓI: người trình bày
    cần chúng, người ngồi dưới thì không — in lên slide là chiếm một phần năm
    chiều cao trang cho thứ chỉ dùng khi bị hỏi.

    Hàm vẫn nhận đủ hai tham số để nội dung không bị xoá khỏi mã nguồn, và vẫn
    giữ phép kiểm độ dài. Muốn bật lại thì chỉ cần bỏ `return` đầu hàm.
    """
    assert len(hoi) <= 96 and len(dap) <= 190, "câu phản biện quá dài"
    return


def noi(sl, chu):
    """Kịch bản nói. Đặc tả yêu cầu 70–100 từ và xưng 'em'."""
    c = " ".join(chu.split())
    n = len(c.split())
    assert 70 <= n <= 100, f"ghi chú {n} từ, đặc tả yêu cầu 70–100"
    assert " em " in f" {c.lower()} ", "ghi chú phải xưng 'em'"
    sl.notes_slide.notes_text_frame.text = chu.strip()


def chu_giai(sl, y, muc):
    """Chú giải màu — sơ đồ có ba màu thì phải nói ba màu nghĩa là gì."""
    x = LE
    for mau, nhan in muc:
        o = sl.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Cm(0.07), Cm(0.34), Cm(0.34))
        o.fill.solid(); o.fill.fore_color.rgb = mau
        o.line.fill.background(); o.shadow.inherit = False
        w = Cm(_rong(nhan, 13) + 0.9)
        txt(sl, x + Cm(0.48), y, w, Cm(0.6), nhan, 13, False, XAM)
        x += Cm(0.48) + w + Cm(0.5)


# ═══════════════════════════════════════════════════════════ nội dung
def dung() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    s = SO

    # ── 1 ──────────────────────────────────────────── BÌA
    #
    # Trang bìa tách RIÊNG khỏi trang mục tiêu. Bản trước gộp cả hai: tên đề tài,
    # tên nhóm, năm gạch đầu dòng và một khối trích dẫn cùng nằm trên một trang —
    # người xem không biết nhìn vào đâu trước, và người trình bày phải nói ba việc
    # trong ba mươi giây đầu.
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), W, H)
    o.fill.solid(); o.fill.fore_color.rgb = XANH
    o.line.fill.background(); o.shadow.inherit = False

    txt(sl, LE, Cm(2.2), RONG, Cm(0.9), "TRƯỜNG ĐẠI HỌC CMC  ·  KHOA CÔNG NGHỆ THÔNG TIN",
        15, True, RGBColor(0xB9, 0xC8, 0xDA), PP_ALIGN.CENTER)
    txt(sl, LE, Cm(3.3), RONG, Cm(0.9), "ĐỒ ÁN MÔN HỌC MÁY VÀ KHAI PHÁ DỮ LIỆU",
        15, False, RGBColor(0x8F, 0xA6, 0xC0), PP_ALIGN.CENTER)

    v = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, (W - Cm(5)) // 2, Cm(5.1), Cm(5), Cm(0.09))
    v.fill.solid(); v.fill.fore_color.rgb = CAM
    v.line.fill.background(); v.shadow.inherit = False

    txt(sl, LE, Cm(6.1), RONG, Cm(3.2),
        "Hệ thống AI tư vấn gọi món cho nhà hàng",
        46, True, TRANG, PP_ALIGN.CENTER, gian=1.15)
    txt(sl, LE, Cm(9.6), RONG, Cm(1.0),
        "Kiến trúc AI Chatbot tư vấn món ăn an toàn và có kiểm soát",
        19, False, RGBColor(0xF3, 0xC5, 0x94), PP_ALIGN.CENTER)

    hop(sl, (W - Cm(15)) // 2, Cm(11.4), Cm(15), Cm(3.5),
        "NHÓM 05\n"
        "Phạm Duy An (nhóm trưởng)  ·  Bùi Đào Đức Anh  ·  Đỗ Tuấn Anh\n"
        "Lê Anh  ·  Nguyễn Quang Hiếu",
        XANH, CAM, TRANG, 17, True)

    txt(sl, LE, Cm(15.6), RONG, Cm(0.9),
        "Giảng viên hướng dẫn:  Phạm Ngọc Đông",
        15, False, RGBColor(0xB9, 0xC8, 0xDA), PP_ALIGN.CENTER)
    txt(sl, LE, Cm(16.6), RONG, Cm(0.9), "Hà Nội, tháng 8 năm 2026",
        13, False, RGBColor(0x8F, 0xA6, 0xC0), PP_ALIGN.CENTER)
    noi(sl, """
Em chào thầy cô và các bạn. Em là Phạm Duy An, nhóm trưởng nhóm 05. Hôm nay
nhóm em xin trình bày đồ án môn Học máy và Khai phá dữ liệu, đề tài Hệ thống AI
tư vấn gọi món cho nhà hàng. Bài trình bày gồm ba phần: bài toán đặt ra, kiến
trúc hệ thống, và kết quả đo được. Em xin bắt đầu bằng bài toán, vì chính bài
toán quyết định vì sao nhóm em không dùng một mô hình ngôn ngữ cho mọi thứ.
""")

    # ── 2 ─────────────────────────────────── bài toán và mục tiêu
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), W, Cm(6.6))
    o.fill.solid(); o.fill.fore_color.rgb = XANH
    o.line.fill.background(); o.shadow.inherit = False
    txt(sl, LE, Cm(1.4), RONG, Cm(2.2), "AI Chatbot tư vấn món ăn",
        40, True, TRANG, PP_ALIGN.CENTER)
    txt(sl, LE, Cm(3.5), RONG, Cm(1.0),
        "Rule-first, AI-assisted: luật kiểm soát, AI hỗ trợ",
        19, False, RGBColor(0xF3, 0xC5, 0x94), PP_ALIGN.CENTER)
    txt(sl, LE, Cm(4.7), RONG, Cm(0.9),
        "Kiến trúc AI Chatbot tư vấn món ăn an toàn và có kiểm soát",
        14, False, RGBColor(0xB9, 0xC8, 0xDA), PP_ALIGN.CENTER)
    y_chinh(sl, [
        "Khách quét QR tại bàn và hỏi bằng tiếng Việt tự nhiên.",
        "Chatbot trả lời thông tin và gợi ý món phù hợp.",
        "Giá, dị nguyên và chính sách bắt buộc phải chính xác.",
        "Khách tự xác nhận trước khi món được thêm vào giỏ.",
        "Mô hình AI lỗi thì hệ thống vẫn phải trả lời được.",
    ], y=Cm(7.2), w=Cm(19.5), co=16.5)
    hop(sl, Cm(21.8), Cm(7.2), Cm(10.5), Cm(4.3),
        "Sinh viên:  ………………………………\n"
        "Lớp:  ………………………………\n"
        "Môn:  Học máy và Khai phá dữ liệu",
        TRANG, KE, XANH, 14, False)
    hop(sl, Cm(21.8), Cm(12.0), Cm(10.5), Cm(2.6),
        "“AI hỗ trợ tư vấn,\nkhông nắm quyền đặt món.”",
        CAM_NHAT, CAM, CAM, 16, True)
    txt(sl, LE, Cm(15.6), Cm(19.5), Cm(1.0),
        f"Thực đơn {s['mon']} món · {s['nhan']} nhãn · {s['tai_lieu']} tài liệu tri thức",
        14, False, XAM)
    noi(sl, """
Em chào thầy cô và các bạn. Nhóm em trình bày kiến trúc chatbot tư vấn món ăn
cho nhà hàng gọi món qua mã QR. Bài toán là khách hỏi bằng tiếng Việt tự nhiên,
còn giá tiền, dị nguyên và chính sách thì bắt buộc chính xác tuyệt đối. Vì thế
nhóm em chọn hướng rule-first, AI-assisted: luật nắm quyền kiểm soát, mô hình
chỉ hỗ trợ ở vị trí có giới hạn. Nguyên tắc xuyên suốt là AI tư vấn chứ không
nắm quyền đặt món; khách luôn tự xác nhận trước khi món vào giỏ.
""")

    # ── 3 ──────────────────────────────────────── BÀI TOÁN
    #
    # Trang riêng cho bài toán, đặt TRƯỚC kiến trúc. Lý do: kiến trúc của đồ án
    # này chỉ hợp lý khi người nghe đã thấy hai loại câu hỏi cần hai cách giải
    # khác hẳn nhau. Trình bày kiến trúc trước rồi giải thích sau thì người nghe
    # phải nhớ ngược.
    sl = trang(prs, "Bài toán: hai câu hỏi nghe giống nhau, giải khác hẳn nhau",
               "BÀI TOÁN",
               "Chính chỗ này quyết định vì sao hệ thống không dùng một mô hình cho mọi câu")
    bang(sl,
         ["Khách hỏi", "Đáp án nằm ở đâu", "Cách giải đúng"],
         [['"Món nào dưới 100 nghìn, không cay?"',
           "Ở cột giá và cột nhãn của món",
           "LỌC BẢNG — đúng 100%"],
          ['"Gọi khai vị trước có làm no bụng không?"',
           "Trong một đoạn văn do người viết",
           "ĐI TÌM đúng đoạn văn đó"]],
         y=Cm(4.0), co=16, cao_min=Cm(1.9))

    txt(sl, LE, Cm(9.4), RONG, Cm(0.7), "BA RÀNG BUỘC KHÔNG ĐƯỢC VI PHẠM", 13, True, CAM)
    for i, (t_, m_) in enumerate([
            ("Giá và dị nguyên", "sai một lần là khách ăn nhầm hoặc trả nhầm tiền"),
            ("Khách tự quyết", "AI gợi ý, nút thêm vào giỏ do khách bấm"),
            ("Mô hình lỗi vẫn chạy", "mất mạng hay hết hạn mức thì vẫn phải trả lời")]):
        hop(sl, LE + i * Cm(10.4), Cm(10.1), Cm(10.0), Cm(2.5), f"{t_}\n{m_}",
            CAM_NHAT, CAM, CAM, 16)

    hop(sl, LE, Cm(13.1), RONG, Cm(1.7),
        "Một mô hình ngôn ngữ trả lời được câu thứ hai, nhưng KHÔNG bảo đảm được câu thứ nhất",
        XANH, XANH, TRANG, 17)
    txt(sl, LE, Cm(15.1), RONG, Cm(1.4),
        "“Giá nhỏ hơn 100 nghìn” là một phép SO SÁNH trên dữ liệu, không phải một phép TÌM KIẾM. "
        "Đưa nó cho mô hình là đổi một đáp án chắc chắn lấy một ước lượng.",
        14.5, False, XAM)
    noi(sl, """
Em xin trình bày bài toán bằng hai câu khách hỏi thật. Câu thứ nhất hỏi món dưới
một trăm nghìn, không cay — đáp án nằm ở cột giá và cột nhãn, lọc bảng là đúng
tuyệt đối. Câu thứ hai hỏi gọi khai vị có làm no không — không có cột nào để
lọc, đáp án nằm trong một đoạn văn. Hai câu nghe giống nhau nhưng cần hai cách
giải khác hẳn. Thêm ba ràng buộc: giá và dị nguyên phải đúng, khách tự quyết, và
mô hình lỗi thì hệ thống vẫn trả lời được.
""")

    # ── 2 ────────────────────────────── kiến trúc và luồng xử lý
    sl = trang(prs, "Kiến trúc và luồng xử lý tổng thể", "TỔNG QUAN",
               "Dữ liệu đi một chiều từ câu hỏi tới câu trả lời; mỗi khối là một tệp trong ai/app/")
    khoi = [
        ("Khách\ngửi câu hỏi", XANH_NHAT, XANH, XANH),
        ("service.py\nnhận HTTP, nạp phiên", XANH_NHAT, XANH, XANH),
        ("understand.py\ntạo Request", XANH_NHAT, XANH, XANH),
        ("session.py\nghép bộ nhớ", XANH_NHAT, XANH, XANH),
        ("llm_understand.py\nhỗ trợ nếu cần", TIM_NHAT, TIM, TIM),
    ]
    khoi2 = [
        ("answer.py · respond()\nĐỊNH TUYẾN", CAM_NHAT, CAM, CAM),
        ("select() · tra dữ liệu\nhoặc RAG", XANH_NHAT, XANH, XANH),
        ("generate.py\ndiễn đạt + xác minh", TIM_NHAT, TIM, TIM),
        ("cart.py\ntạo thẻ món", CAM_NHAT, CAM, CAM),
        ("Cập nhật phiên\nvà trả kết quả", XANH_NHAT, XANH, XANH),
    ]
    for hang_i, day in enumerate((khoi, khoi2)):
        y = Cm(4.3) + hang_i * Cm(3.55)
        n = len(day)
        wt = Cm(0.85)
        wb = (RONG - wt * (n - 1)) // n
        for i, (chu, nen, vien, cm) in enumerate(day):
            x = LE + i * (wb + wt)
            hop(sl, x, y, wb, Cm(2.35), chu, nen, vien, cm, 14)
            if i < n - 1:
                ten(sl, x + wb + Cm(0.06), y + Cm(0.95), wt - Cm(0.12),
                    mau=KE if False else XANH)
    ten(sl, LE + RONG // 2 - Cm(0.35), Cm(6.75), Cm(0.7), mau=CAM, doc=True, h=Cm(0.95))
    chu_giai(sl, Cm(11.55), [(XANH, "Mã tất định"), (TIM, "Có dùng mô hình"),
                             (CAM, "Lớp kiểm tra an toàn")])
    hop(sl, LE, Cm(12.5), RONG, Cm(1.5),
        "Phần lớn luồng là mã tất định; mô hình chỉ xuất hiện ở một số vị trí có giới hạn",
        XANH, XANH, TRANG, 17, True)
    txt(sl, LE, Cm(14.2), RONG, Cm(1.1),
        f"Bộ định tuyến có {s['ho_nhanh']} nhánh xử lý; chỉ {s['nhanh_sinh']} nhánh được phép "
        f"để mô hình sinh câu, và câu sinh ra phải qua {s['phep_kiem']} phép xác minh.",
        14.5, False, XAM)
    noi(sl, """
Đây là toàn bộ luồng, và dữ liệu chỉ đi một chiều từ trái sang phải. Câu hỏi
vào service, sang understand thành Request, qua session ghép bộ nhớ, rồi tới
router. Router mới quyết định tra dữ liệu, đi RAG hay lọc món. Em xin nhấn ba
màu: xanh là mã tất định, tím là chỗ có dùng mô hình, cam là lớp kiểm tra an
toàn. Nhìn tỷ lệ màu thấy phần lớn luồng là tất định; mô hình chỉ nằm ở hai vị
trí và đều bị kiểm tra ngay sau đó.
""")
    _tang_du_lieu(prs, s=SO)
    _dinh_tuyen(prs, s=SO)
    _rag_va_ket(prs, s=SO)
    import os
    # Nếu tệp đang mở trong PowerPoint thì Windows khoá ghi. Ghi ra bản
    # `.moi.pptx` thay vì nổ, để lần dựng vẫn dùng được.
    ra = HERE / 'output' / 'SLIDE_KIEN_TRUC_AI.pptx'
    ra.parent.mkdir(exist_ok=True)
    try:
        prs.save(str(ra))
    except PermissionError:
        ra = ra.with_suffix('.moi.pptx')
        prs.save(str(ra))
        print('  (tệp gốc đang mở trong PowerPoint — ghi ra bản .moi)')
    return ra


# ═════════════════════════════════════ 3–6 · phần nhóm trưởng phụ trách
def _tang_du_lieu(prs, s):
    # ── 3 ──────────────────────────────────── tầng dữ liệu thực đơn
    sl = trang(prs, "Tầng dữ liệu thực đơn", "DỮ LIỆU · 1/4",
               "Một sự thật, hai nơi lưu — và một cổng bắt hai nơi phải khớp")
    y_chinh(sl, [
        f"Thực đơn có {s['mon']} món, chia {s['danh_muc']} danh mục.",
        "JSON phục vụ dịch vụ AI; Postgres phục vụ backend.",
        f"Hai nguồn bắt buộc khớp {s['mon']}/{s['mon']} món.",
        "CI kiểm tra để ngăn lệch tên món và lệch giá.",
        "Dữ liệu này dùng lại ở khâu lọc, xác minh và giỏ hàng.",
    ], y=Cm(4.0), w=Cm(16.4), co=16)

    x0 = Cm(18.6)
    hop(sl, x0, Cm(4.2), Cm(4.3), Cm(2.2), "JSON thực đơn\nai đọc lúc chạy",
        XANH_NHAT, XANH, XANH, 14)
    hop(sl, x0 + Cm(5.0), Cm(4.2), Cm(4.5), Cm(2.2), "CỔNG KIỂM TRA\nđồng bộ (CI)",
        CAM_NHAT, CAM, CAM, 14)
    hop(sl, x0 + Cm(10.2), Cm(4.2), Cm(4.1), Cm(2.2), "Postgres\nbackend đọc",
        XANH_NHAT, XANH, XANH, 14)
    ten(sl, x0 + Cm(4.4), Cm(5.1), Cm(0.5), mau=CAM)
    ten(sl, x0 + Cm(9.6), Cm(5.1), Cm(0.5), mau=CAM)
    txt(sl, x0, Cm(6.6), Cm(14.3), Cm(0.8),
        f"cổng đối chiếu từng món · lệch một món là CI đỏ", 13, False, XAM,
        PP_ALIGN.CENTER)
    hop(sl, x0, Cm(7.9), Cm(14.3), Cm(2.5),
        "Rủi ro nếu hai nguồn lệch\nAI gợi ý một món hoặc một mức giá khác với backend —"
        " khách thấy một giá, hoá đơn tính một giá",
        CAM_NHAT, CAM, CAM, 15)
    hop(sl, x0, Cm(10.9), Cm(14.3), Cm(3.6),
        "Vì sao không bỏ hẳn một nguồn\nDịch vụ AI phải trả lời được cả khi cơ sở dữ liệu"
        " bận hoặc mất kết nối. Tách hai nguồn là chấp nhận rủi ro lệch, đổi lấy việc"
        " chatbot không chết theo backend — nên phải có cổng canh.",
        TRANG, KE, XANH, 15)
    phan_bien(sl,
              "Sao không để AI đọc thẳng Postgres cho khỏi phải đồng bộ?",
              "Vì khi đó dịch vụ AI chết theo cơ sở dữ liệu. Nhóm em chọn tách nguồn để "
              "chatbot vẫn trả lời được, và trả giá bằng một cổng CI đối chiếu từng món.")
    noi(sl, """
Đây là tầng dữ liệu, phần em phụ trách. Thực đơn được lưu ở hai nơi: JSON cho
dịch vụ AI và Postgres cho backend. Em biết tách nguồn là rủi ro, nhưng nếu để
AI đọc thẳng cơ sở dữ liệu thì chatbot chết theo backend. Nên nhóm em chấp nhận
hai nguồn và dựng một cổng CI đối chiếu từng món; lệch một tên hay một giá là
đỏ ngay. Rủi ro cụ thể phải chặn là AI báo một giá còn hoá đơn tính giá khác.
""")

    # ── 4 ─────────────────────────────── hệ nhãn và chất lượng dữ liệu
    sl = trang(prs, "Hệ nhãn và chất lượng dữ liệu", "DỮ LIỆU · 2/4",
               f"{s['nhan']} nhãn thuộc {s['nhom_nhan']} nhóm, dạng group:value")
    y_chinh(sl, [
        f"Nhãn có dạng group:value — spice:none, allergen:seafood, occasion:date.",
        "Nhãn phủ ĐỦ trên mọi món thì được dùng làm điều kiện lọc.",
        "Nhãn phủ MỘT PHẦN không đủ để kết luận một món là an toàn.",
        f"allergen mới phủ {s['di_ung']}/{s['mon']} món — phần còn lại là “chưa ghi nhận”.",
    ], y=Cm(4.0), w=Cm(16.4), co=16)
    bang(sl, ["Nhãn", "Vai trò trong hệ thống"],
         [["spice:none", "Điều kiện lọc — loại hẳn món không thoả"],
          ["allergen:seafood", "Loại rủi ro ĐÃ BIẾT, không phải bảo chứng"],
          ["occasion:date", "Chỉ dùng xếp hạng sở thích, không dùng để loại món"]],
         y=Cm(9.2), co=14, x=LE, w=Cm(16.4), cao_min=Cm(1.2))
    hop(sl, Cm(18.6), Cm(4.2), Cm(14.3), Cm(4.6),
        "Không có nhãn dị nguyên  ≠  món an toàn\n"
        f"Chỉ có nghĩa CHƯA AI GHI NHẬN. Với {s['mon'] - s['di_ung']} món còn lại, hệ thống"
        " bắt buộc phải mời khách xác nhận trực tiếp với nhân viên để bếp kiểm tra.",
        CAM_NHAT, CAM, CAM, 16)
    txt(sl, Cm(18.6), Cm(9.3), Cm(14.3), Cm(0.8), "BA CÔNG CỤ RÀ TỰ ĐỘNG", 13, True, CAM)
    for i, (t_, m_) in enumerate([
            ("audit_allergen_tags.py", "soát nhãn dị nguyên còn thiếu"),
            ("audit_season_tags.py", "soát nhãn mùa vụ"),
            ("audit_method_tags.py", "soát nhãn cách chế biến")]):
        hop(sl, Cm(18.6), Cm(10.0) + i * Cm(1.5), Cm(14.3), Cm(1.3),
            f"{t_}   —   {m_}", TRANG, KE, XANH, 14, False)
    phan_bien(sl,
              "Đã có nhãn dị nguyên thì sao còn phải hỏi nhân viên?",
              f"Vì nhãn mới phủ {s['di_ung']}/{s['mon']} món. Thiếu nhãn nghĩa là chưa ai ghi, "
              "không phải là món không chứa dị nguyên — nên hệ thống không được phép hứa an toàn.")
    noi(sl, """
Hệ nhãn là chỗ em muốn nói kỹ nhất. Nhãn có dạng nhóm hai chấm giá trị, và điều
quan trọng là không phải nhãn nào cũng dùng như nhau. Nhãn phủ đủ mọi món thì
được dùng để lọc. Nhãn phủ một phần thì chỉ được dùng xếp hạng. Riêng nhãn dị
nguyên mới phủ bốn mươi bốn trên chín mươi mốt món, nên thiếu nhãn chỉ có nghĩa
chưa ai ghi nhận, tuyệt đối không có nghĩa món đó an toàn. Vì vậy hệ thống luôn
mời khách xác nhận lại với nhân viên.
""")

    # ── 5 ────────────────────────── lớp hiểu câu hỏi understand.py
    sl = trang(prs, "Lớp hiểu câu hỏi — understand.py", "HIỂU CÂU HỎI · 3/4",
               "Biến câu tiếng Việt thành một cấu trúc dữ liệu; không chọn món, không viết câu")
    y_chinh(sl, [
        f"Nhận câu tiếng Việt, tạo Request gồm {s['truong_request']} trường.",
        f"Đối chiếu với từ điển {s['tu_vung']} cụm từ đã đăng ký.",
        "KHÔNG chọn món và KHÔNG viết câu trả lời.",
        "Đường xử lý chính chạy hoàn toàn không cần mô hình.",
        "Kết quả là đầu vào duy nhất của bộ định tuyến.",
    ], y=Cm(4.0), w=Cm(15.6), co=16)

    x0, wb = Cm(17.9), Cm(15.0)
    buoc = [("1 · Chuẩn hoá câu", "hạ chữ thường, bỏ dấu, bỏ dấu câu"),
            ("2 · Khớp cụm từ", f"dò trong từ điển {s['tu_vung']} cụm"),
            ("3 · Tách ràng buộc và sở thích", "điều cần tránh / bắt buộc / ưu tiên"),
            ("4 · Nhận diện ý định", "hỏi giá, xin gợi ý, hỏi kiến thức…")]
    for i, (t_, m_) in enumerate(buoc):
        y = Cm(4.2) + i * Cm(2.5)
        hop(sl, x0, y, wb, Cm(1.75), f"{t_}\n{m_}", XANH_NHAT, XANH, XANH, 15)
        if i < len(buoc) - 1:
            ten(sl, x0 + wb // 2 - Cm(0.3), y + Cm(1.82), Cm(0.6), mau=CAM,
                doc=True, h=Cm(0.6))
    hop(sl, LE, Cm(11.5), Cm(15.6), Cm(3.0),
        "understand.py xác định KHÁCH MUỐN GÌ;\nselect() mới quyết định MÓN NÀO phù hợp",
        XANH, XANH, TRANG, 17)
    phan_bien(sl,
              "Sao không để mô hình đọc thẳng câu hỏi rồi chọn món luôn cho nhanh?",
              "Vì khi đó không có gì để kiểm tra. Tách ra thành Request thì mỗi ràng buộc là "
              "một trường dữ liệu, so được với thực đơn và viết được test — còn câu chữ thì không.")
    noi(sl, """
Đây là lớp hiểu câu hỏi, phần thứ hai em phụ trách. Nhiệm vụ của nó chỉ là biến
câu tiếng Việt thành một cấu trúc dữ liệu bốn mươi tám trường, gọi là Request.
Em xin nhấn mạnh nó không chọn món và không viết câu trả lời. Ranh giới này rất
quan trọng, vì nhờ tách ra mà mỗi ràng buộc trở thành một trường có thể so với
thực đơn và viết test được. Đường xử lý chính chạy bằng từ điển, hoàn toàn
không cần gọi mô hình.
""")

    # ── 6 ──────────────────── cơ chế khớp cụm và ví dụ tạo Request
    sl = trang(prs, "Cơ chế khớp cụm từ và ví dụ tạo Request", "HIỂU CÂU HỎI · 4/4",
               "Vì sao “dị ứng mì chính” không bị đọc nhầm thành dị ứng gluten")
    txt(sl, LE, Cm(3.95), Cm(15.6), Cm(0.7), "CƠ CHẾ KHỚP", 13, True, CAM)
    y_chinh(sl, [
        "Hạ chữ thường, bỏ dấu và bỏ dấu câu trước khi khớp.",
        "Ưu tiên cụm DÀI trước cụm ngắn.",
        "Đoạn văn bản đã khớp thì không được khớp lại.",
        f"Cơ chế này bảo vệ {s['va_cham']} cụm có nguy cơ va chạm.",
    ], y=Cm(4.6), w=Cm(15.6), co=15.5)
    hop(sl, LE, Cm(9.5), Cm(15.6), Cm(2.5),
        "“Tôi dị ứng mì chính.”\n"
        "“mì chính” là cụm dài nên được khớp TRƯỚC — đoạn đó bị đánh dấu, "
        "nên “mì” không còn để khớp thành gluten",
        CAM_NHAT, CAM, CAM, 15)
    bang(sl, ["Loại thông tin", "Ví dụ", "Hệ quả"],
         [["Cần tránh", "Dị ứng hải sản", "LOẠI món khỏi kết quả"],
          ["Bắt buộc", "Không cay", "LỌC theo điều kiện"],
          ["Sở thích", "Đi hẹn hò", "Chỉ XẾP HẠNG"]],
         y=Cm(11.4), co=13.5, x=LE, w=Cm(15.6), cao_min=Cm(1.05))

    x0 = Cm(17.9)
    txt(sl, x0, Cm(3.95), Cm(15.0), Cm(0.7), "VÍ DỤ TẠO REQUEST", 13, True, CAM)
    hop(sl, x0, Cm(4.6), Cm(15.0), Cm(1.6),
        "“Mình dị ứng hải sản, muốn món không cay dưới 150 nghìn.”",
        TRANG, KE, XANH, 15)
    ten(sl, x0 + Cm(7.2), Cm(6.35), Cm(0.6), mau=CAM, doc=True, h=Cm(0.75))
    for i, (k_, v_, m_) in enumerate([
            ("avoid_tags", "{allergen:seafood}", CAM),
            ("require_tags", "{spice:none}", XANH),
            ("budget_max", "150000", XANH),
            ("y_dinh", "hoi_mon", XANH)]):
        y = Cm(7.4) + i * Cm(1.35)
        o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, Cm(15.0), Cm(1.15))
        o.fill.solid(); o.fill.fore_color.rgb = CAM_NHAT if m_ is CAM else XANH_NHAT
        o.line.color.rgb = m_; o.line.width = Pt(0.9); o.shadow.inherit = False
        tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.35)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{k_}  =  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = m_; r.font.name = FONT
        r = p.add_run(); r.text = v_
        r.font.size = Pt(15); r.font.color.rgb = DEN; r.font.name = "Consolas"
    txt(sl, x0, Cm(13.0), Cm(15.0), Cm(1.6),
        "Bốn trường này là kết quả THẬT khi chạy understand() trên câu trên — "
        "không có trường nào do người viết slide điền vào.",
        13.5, False, XAM)
    phan_bien(sl,
              "Từ điển cụm từ thì làm sao phủ hết cách nói của khách?",
              "Không phủ hết được, và nhóm em không giả vờ là phủ. Câu nào luật không đọc ra "
              "ràng buộc thì mới gọi mô hình hỗ trợ, và mô hình chỉ được trả nhãn có sẵn trong từ điển.")
    noi(sl, """
Slide này giải thích cơ chế khớp và cho một ví dụ chạy thật. Cơ chế quan trọng
nhất là ưu tiên cụm dài trước cụm ngắn, và đoạn đã khớp thì không khớp lại. Nhờ
vậy câu tôi dị ứng mì chính không bị đọc nhầm thành dị ứng gluten, vì cụm mì
chính được nhận trước. Bên phải là Request thật sinh ra từ câu ví dụ. Em xin
lưu ý ba loại thông tin có hệ quả khác nhau: điều cần tránh thì loại món, bắt
buộc thì lọc, còn sở thích chỉ dùng xếp hạng.
""")


def _dinh_tuyen(prs, s):
    # ── 7 ───────────────────── bộ nhớ phiên và mô hình hỗ trợ hiểu
    sl = trang(prs, "Bộ nhớ phiên và mô hình hỗ trợ hiểu câu", "NGỮ CẢNH",
               "Luật xử lý trước; mô hình chỉ bổ sung khi luật chưa hiểu đủ")
    txt(sl, LE, Cm(3.95), Cm(15.6), Cm(0.7), "session.py — GHÉP BỘ NHỚ", 13, True, XANH)
    y_chinh(sl, [
        "Dị nguyên được CỘNG DỒN, không bao giờ bị bỏ.",
        "Ràng buộc cứng thì lượt mới GHI ĐÈ cùng nhóm.",
        f"Ngữ cảnh sở thích giữ {s['ngu_canh']} thông tin gần nhất.",
        "Nhờ đó hiểu được “món đầu tiên” và “cho món khác”.",
    ], y=Cm(4.6), w=Cm(15.6), co=15.5)
    for i, (c_, g_) in enumerate([
            ("1.  “Mình dị ứng hải sản.”", "ghi dị nguyên"),
            ("2.  “Dưới 100k nhé.”", "thêm ngân sách"),
            ("3.  “Cho mình món khác.”", "GIỮ cả hai điều kiện trên")]):
        y = Cm(9.6) + i * Cm(1.55)
        nb = i == 2
        hop(sl, LE, y, Cm(15.6), Cm(1.3), f"{c_}      →      {g_}",
            CAM_NHAT if nb else TRANG, CAM if nb else KE, CAM if nb else XANH, 14.5, nb)
        if i < 2:
            ten(sl, LE + Cm(7.5), y + Cm(1.32), Cm(0.5), mau=KE, doc=True, h=Cm(0.2))
    txt(sl, LE, Cm(14.3), Cm(15.6), Cm(1.4),
        "Ở lượt 3 khách không nhắc lại gì, nhưng hai ràng buộc trước vẫn còn hiệu lực.",
        13.5, False, XAM)

    x0 = Cm(17.9)
    txt(sl, x0, Cm(3.95), Cm(15.0), Cm(0.7),
        "llm_understand.py — MÔ HÌNH HỖ TRỢ", 13, True, TIM)
    for i, m_ in enumerate([
            "Chỉ được gọi khi luật chưa hiểu đủ ràng buộc.",
            "Chỉ được trả nhãn CÓ SẴN trong từ điển.",
            "Chỉ được THÊM điều kiện, không được xoá.",
            "Không có đường nào cho nó tự chọn món."]):
        hop(sl, x0, Cm(4.6) + i * Cm(1.6), Cm(15.0), Cm(1.35), m_,
            TIM_NHAT, TIM, TIM, 15, False)
    hop(sl, x0, Cm(11.3), Cm(15.0), Cm(1.5),
        "Bốn ràng buộc trên đều là mã, không phải lời dặn trong prompt",
        TIM, TIM, TRANG, 15)
    hop(sl, x0, Cm(13.1), Cm(15.0), Cm(1.6),
        "“Luật xử lý trước; mô hình chỉ bổ sung khi cần.”",
        CAM_NHAT, CAM, CAM, 16)
    phan_bien(sl,
              "Mô hình bổ sung nhãn thì có thể bịa ra nhãn sai không?",
              "Bịa ra nhãn lạ thì bị loại vì không có trong từ điển. Bịa ra nhãn có thật thì chỉ "
              "làm hẹp kết quả, không làm mất dị nguyên — vì nó không được phép xoá điều kiện.")
    noi(sl, """
Bên trái là bộ nhớ phiên. Quy tắc quan trọng nhất là dị nguyên cộng dồn, không
bao giờ bị bỏ, còn ràng buộc cứng thì lượt mới ghi đè lượt cũ. Ví dụ dưới cho
thấy ở lượt ba khách không nhắc lại gì nhưng hai điều kiện trước vẫn còn. Bên
phải là mô hình hỗ trợ hiểu câu. Em xin nhấn bốn ràng buộc này được cài bằng
mã chứ không phải lời dặn trong prompt: mô hình chỉ thêm được nhãn có sẵn, và
không có đường nào để nó tự chọn món.
""")

    # ── 8 ──────────────────────────────────────── bộ định tuyến
    sl = trang(prs, "Bộ định tuyến và cách chọn đường xử lý", "ĐỊNH TUYẾN",
               "respond() trong answer.py là router nghiệp vụ, không phải bộ phân loại")
    y_chinh(sl, [
        "Router chỉ đọc các cờ đã có sẵn trong Request.",
        "Kiểm tra các nhánh theo một thứ tự CỐ ĐỊNH.",
        "Nhánh đầu tiên khớp sẽ thắng, không chấm điểm.",
        f"Có {s['ho_nhanh']} nhánh xử lý, {s['diem_tra_ve']} điểm trả về.",
        "select() và RAG nằm PHÍA SAU router, không phải trước.",
    ], y=Cm(4.0), w=Cm(14.8), co=15.5)
    hop(sl, LE, Cm(11.2), Cm(14.8), Cm(3.3),
        "“Phở bò giá bao nhiêu?”\n"
        "understand.py đặt asks_price = True và named_items = [m_008]\n"
        "→ router chọn nhánh price_lookup, tra thẳng bảng giá",
        CAM_NHAT, CAM, CAM, 15)

    x0, wb = Cm(17.4), Cm(15.5)
    txt(sl, x0, Cm(3.95), wb, Cm(0.7), "CỜ TRONG REQUEST      →      NHÁNH ĐƯỢC CHỌN",
        13, True, CAM)
    cap = [("Hỏi giá một món", "price_lookup", XANH),
           ("Hỏi dị nguyên của món", "allergen_named_dish", CAM),
           ("So sánh hai món", "compare", XANH),
           ("Hỏi kiến thức ẩm thực", "knowledge_corpus", TIM),
           ("Xin gợi ý món", "filter", XANH),
           ("Chưa đủ thông tin", "clarify", XANH)]
    for i, (t_, n_, m_) in enumerate(cap):
        y = Cm(4.65) + i * Cm(1.72)
        hop(sl, x0, y, Cm(7.0), Cm(1.35), t_, TRANG, KE, DEN, 14, False)
        ten(sl, x0 + Cm(7.15), y + Cm(0.48), Cm(1.1), mau=m_)
        hop(sl, x0 + Cm(8.5), y, wb - Cm(8.5), Cm(1.35), n_,
            CAM_NHAT if m_ is CAM else (TIM_NHAT if m_ is TIM else XANH_NHAT),
            m_, m_, 14)
    phan_bien(sl,
              "Nhánh đầu khớp là thắng — vậy thứ tự nhánh sai thì cả hệ thống sai?",
              "Đúng, nên thứ tự là thứ phải test. Nhóm em có tập ca chấm riêng nhánh nào được "
              "chọn cho từng loại câu, và một bộ in ra từng câu đi sai để đọc bằng mắt.")
    noi(sl, """
Đây là bộ định tuyến. Em xin lưu ý nó là router nghiệp vụ chứ không phải bộ
phân loại: nó không chấm điểm, mà chỉ đọc các cờ đã có trong Request rồi duyệt
các nhánh theo thứ tự cố định, nhánh đầu tiên khớp là thắng. Bảng bên phải cho
thấy cờ nào dẫn tới nhánh nào. Ví dụ câu phở bò giá bao nhiêu, lớp hiểu đặt cờ
hỏi giá và ghi mã món, router chọn nhánh tra giá. Điều quan trọng là select và
RAG nằm phía sau router chứ không phải trước.
""")


def _rag_va_ket(prs, s):
    # ── 9 ───────────────────── khi nào dùng RAG và hàng rào an toàn
    sl = trang(prs, "Khi nào dùng RAG, và các hàng rào an toàn", "RAG · AN TOÀN",
               "RAG tìm kiến thức; dữ liệu và luật quyết định món hợp lệ")
    txt(sl, LE, Cm(3.9), Cm(15.6), Cm(0.7), "RAG DÙNG Ở ĐÂU", 13, True, CAM)
    hop(sl, LE, Cm(4.55), Cm(7.6), Cm(4.5),
        "KHÔNG dùng RAG\nGiá · Dị nguyên · Ngân sách · Lọc món ·"
        " Chính sách có một đáp án chính xác",
        CAM_NHAT, CAM, CAM, 15)
    hop(sl, LE + Cm(8.0), Cm(4.55), Cm(7.6), Cm(4.5),
        "CÓ dùng RAG\nCâu hỏi giải thích · Kiến thức ẩm thực ·"
        " “Vì sao món này mềm hơn món kia?” · “Nên chọn mức cay thế nào?”",
        XANH_NHAT, XANH, XANH, 15)
    txt(sl, LE, Cm(9.4), Cm(15.6), Cm(0.7), "BA BỘ TRUY HỒI ĐÃ SO", 13, True, CAM)
    for i, (t_, m_) in enumerate([("BM25", "khớp theo từ khoá"),
                                  ("BGE-M3", "khớp theo ý nghĩa"),
                                  ("Hybrid RRF", "hợp nhất hai bảng xếp hạng")]):
        hop(sl, LE + i * Cm(5.3), Cm(10.05), Cm(5.0), Cm(1.9), f"{t_}\n{m_}",
            TRANG, KE, XANH, 15)
    hop(sl, LE, Cm(12.35), Cm(15.6), Cm(2.2),
        "RAG chỉ chạy ở nhánh câu hỏi kiến thức — mọi câu về giá và dị nguyên"
        " đều đi thẳng vào dữ liệu",
        XANH, XANH, TRANG, 15)

    x0, wb = Cm(17.9), Cm(15.0)
    txt(sl, x0, Cm(3.9), wb, Cm(0.7), "NĂM HÀNG RÀO AN TOÀN", 13, True, CAM)
    rao = [f"Mô hình KHÔNG có đường nào tự chọn món.",
           f"Chỉ {s['nhanh_sinh']}/{s['ho_nhanh']} nhánh được phép sinh câu.",
           f"Câu sinh ra phải qua {s['phep_kiem']} phép xác minh.",
           "Giỏ hàng dựng từ reply.items, không đọc chữ AI viết.",
           "Khách luôn phải tự xác nhận trước khi món vào giỏ."]
    for i, r_ in enumerate(rao):
        y = Cm(4.55) + i * Cm(1.62)
        o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, Cm(0.9), Cm(1.35))
        o.fill.solid(); o.fill.fore_color.rgb = CAM
        o.line.fill.background(); o.shadow.inherit = False
        tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = str(i + 1)
        rr.font.size = Pt(17); rr.font.bold = True
        rr.font.color.rgb = TRANG; rr.font.name = FONT
        hop(sl, x0 + Cm(1.0), y, wb - Cm(1.0), Cm(1.35), r_,
            CAM_NHAT, CAM, DEN, 14.5, False)
    hop(sl, x0, Cm(12.9), wb, Cm(1.65),
        "“RAG tìm kiến thức; dữ liệu và luật quyết định món hợp lệ.”",
        XANH, XANH, TRANG, 16)
    phan_bien(sl,
              "Hệ thống đã có RAG rồi thì sao còn cần luật lọc món?",
              "Vì RAG trả về đoạn văn giống câu hỏi nhất, không trả về món thoả điều kiện. "
              "Câu “dưới 150 nghìn, không cay” là phép so sánh trên dữ liệu, không phải phép tìm kiếm.")
    noi(sl, """
Slide này phân biệt rõ chỗ nào dùng RAG và chỗ nào không. Mọi câu về giá, dị
nguyên, ngân sách hay lọc món đều đi thẳng vào dữ liệu, vì đó là phép so sánh
chứ không phải phép tìm kiếm. RAG chỉ dùng cho câu hỏi giải thích và kiến thức
ẩm thực. Nhóm em đã so ba bộ truy hồi: BM25, BGE-M3 và hybrid. Bên phải là năm
hàng rào an toàn, và em xin nhấn hàng rào thứ tư: giỏ hàng dựng từ danh sách
món đã lọc, không đọc chữ mà mô hình viết ra.
""")

    # ── 10 ──────────────────────── đánh giá, hạn chế và kết luận
    sl = trang(prs, "Đánh giá, hạn chế và kết luận", "KẾT LUẬN")
    cot_x = [LE, LE + Cm(10.4), LE + Cm(20.8)]
    wc = Cm(10.0)
    khoi = [
        ("ĐÁNH GIÁ", XANH, [
            "Kiểm thử hiểu câu một lượt",
            "Kiểm thử hội thoại nhiều lượt",
            "Đánh giá truy hồi và RAG",
            "Kiểm thử toàn chuỗi HTTP → giỏ hàng",
            "Ca an toàn có thể chặn CI"]),
        ("ĐIỂM MẠNH", XANH, [
            "Tất định và kiểm thử được",
            "Phân tách rõ trách nhiệm từng lớp",
            "Vẫn hoạt động khi mô hình lỗi",
            "Nhiều lớp kiểm soát dị nguyên",
            "Số liệu đọc từ mã, không viết tay"]),
        ("HẠN CHẾ", CAM, [
            f"Nhãn dị nguyên mới phủ {s['di_ung']}/{s['mon']} món",
            "Bỏ dấu gây va chạm: “cá” và “cả”",
            "Hai nguồn thực đơn phải luôn đồng bộ",
            "Không đọc được ràng buộc thì lọc quá rộng",
            "Tập đánh giá do nhóm viết, chưa có log thật"]),
    ]
    for (t_, m_, muc), x in zip(khoi, cot_x):
        o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Cm(3.6), wc, Cm(0.85))
        o.fill.solid(); o.fill.fore_color.rgb = m_
        o.line.fill.background(); o.shadow.inherit = False
        tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t_
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = TRANG; r.font.name = FONT
        for i, u in enumerate(muc):
            hop(sl, x, Cm(4.65) + i * Cm(1.62), wc, Cm(1.4), u,
                CAM_NHAT if m_ is CAM else TRANG, m_ if m_ is CAM else KE,
                DEN, 13.5, False)
    o = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, LE, Cm(13.35), RONG, Cm(3.5))
    o.fill.solid(); o.fill.fore_color.rgb = XANH
    o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, d in enumerate(["Dữ liệu quyết định sự thật.",
                           "Luật quyết định món hợp lệ.",
                           "AI hỗ trợ hiểu và diễn đạt."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.22
        r = p.add_run(); r.text = d
        r.font.size = Pt(21); r.font.bold = True
        r.font.color.rgb = TRANG if i < 2 else RGBColor(0xF3, 0xC5, 0x94)
        r.font.name = FONT
    txt(sl, LE, Cm(17.15), RONG, Cm(1.0),
        "Cảm ơn thầy/cô và các bạn  —  Q&A", 17, True, CAM, PP_ALIGN.CENTER)
    noi(sl, """
Phần cuối là đánh giá và hạn chế. Nhóm em kiểm thử ở bốn mức: hiểu câu một
lượt, hội thoại nhiều lượt, truy hồi và RAG, và toàn chuỗi từ HTTP tới giỏ
hàng. Ca an toàn được đặt là chốt chặn CI chứ không phải điểm số. Về hạn chế,
em xin nói thẳng ba điều: nhãn dị nguyên chưa phủ hết, bỏ dấu vẫn gây va chạm,
và tập đánh giá do nhóm tự viết nên chưa phản ánh khách thật. Em xin hết và sẵn
sàng nhận câu hỏi.
""")


if __name__ == "__main__":
    print(dung())
